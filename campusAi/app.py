from flask import Flask, render_template, request, jsonify
from werkzeug.utils import secure_filename
from docx import Document
import os
from pptx import Presentation
from pdfminer.high_level import extract_text
from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer
from sumy.summarizers.lsa import LsaSummarizer

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max file size

# Ensure upload folder exists
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# Allowed file extensions
ALLOWED_EXTENSIONS = {'pdf', 'pptx'}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_docx(file_path):
    """Extracts table data from a Word document."""
    doc = Document(file_path)
    tables_data = []
    for table in doc.tables:
        table_content = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        tables_data.append(table_content)
    return tables_data

def format_timetable(table_data):
    """Formats extracted timetable table into a dictionary."""
    timetable = {}
    headers = table_data[0][1:]
    for row in table_data[1:]:
        if row:
            day = row[0]
            classes = row[1:]
            timetable[day] = {headers[i]: classes[i] for i in range(len(classes))}
    return timetable

def extract_text_from_ppt(ppt_path):
    presentation = Presentation(ppt_path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
    return text.strip()

def extract_text_from_pdf(pdf_path):
    return extract_text(pdf_path).strip()

def summarize_text(text, num_sentences=5):
    parser = PlaintextParser.from_string(text, Tokenizer("english"))
    summarizer = LsaSummarizer()
    summary = summarizer(parser.document, num_sentences)
    return " ".join(str(sentence) for sentence in summary)

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400
    
    if not file.filename.endswith('.docx'):
        return jsonify({'error': 'Only .docx files are allowed'}), 400
    
    try:
        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(filepath)
        
        tables_data = extract_text_from_docx(filepath)
        if not tables_data:
            return jsonify({'error': 'No tables found in document'}), 400
        
        timetable = format_timetable(tables_data[0])
        os.remove(filepath)  # Clean up uploaded file
        
        return jsonify({'timetable': timetable})
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/study-groups')
def study_groups():
    return render_template('study_groups.html')

@app.route('/reminders')
def reminders():
    return render_template('reminders.html')

@app.route('/attendance')
def attendance():
    return render_template('attendance.html')

@app.route('/login')
def login():
    return render_template('login.html')

@app.route('/summarizer')
def summarizer():
    return render_template('summarizer.html')

@app.route('/summarize', methods=['POST'])
def summarize():
    if 'file' not in request.files:
        return jsonify({'success': False, 'error': 'No file uploaded'})
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'success': False, 'error': 'No file selected'})
    
    if not allowed_file(file.filename):
        return jsonify({'success': False, 'error': 'Invalid file type'})
    
    try:
        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(filepath)
        
        # Extract text based on file type
        if filename.endswith('.pptx'):
            text = extract_text_from_ppt(filepath)
        else:  # PDF
            text = extract_text_from_pdf(filepath)
        
        # Generate summary
        summary = summarize_text(text)
        
        # Clean up uploaded file
        os.remove(filepath)
        
        return jsonify({
            'success': True,
            'summary': summary
        })
        
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e)
        })

if __name__ == '__main__':
    app.run(debug=True, port=5001) 